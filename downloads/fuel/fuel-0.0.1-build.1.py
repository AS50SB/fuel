import os
import sys
import zipfile
import tempfile
import shutil
from io import BytesIO
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
import tkinter as tk
from tkinter import ttk, messagebox, filedialog

class MultiJarPPTMerger:
    def __init__(self, root):
        self.root = root
        self.root.title("Fuel Loader-0.0.1-build.1")
        self.root.geometry("800x600")
        
        # 变量初始化
        self.target_pptx_path = None
        self.mods_dir = None
        self.selected_jars = []
        self.temp_dir = tempfile.mkdtemp()
        
        # 创建UI
        self._create_ui()
        
        # 启动时自动搜索mods文件夹
        self._auto_search_mods()

    def _create_ui(self):
        """创建用户界面"""
        main_frame = ttk.Frame(self.root, padding="20")
        main_frame.pack(fill=tk.BOTH, expand=True)
        
        # 目标PPTX文件选择
        ttk.Label(main_frame, text="SFC路径:").grid(row=0, column=0, padx=5, pady=10, sticky=tk.W)
        self.target_pptx_var = tk.StringVar(value="未选择")
        ttk.Entry(main_frame, textvariable=self.target_pptx_var, width=60, state="readonly").grid(row=0, column=1, padx=5, pady=10)
        ttk.Button(main_frame, text="浏览", command=self._select_target_pptx).grid(row=0, column=2, padx=5, pady=10)
        
        # mods文件夹路径
        ttk.Label(main_frame, text="mods文件夹:").grid(row=1, column=0, padx=5, pady=10, sticky=tk.W)
        self.mods_dir_var = tk.StringVar(value="搜索中...")
        ttk.Entry(main_frame, textvariable=self.mods_dir_var, width=60, state="readonly").grid(row=1, column=1, padx=5, pady=10)
        ttk.Button(main_frame, text="浏览", command=self._select_mods_dir).grid(row=1, column=2, padx=5, pady=10)
        
        # JAR文件列表
        ttk.Label(main_frame, text="mod列表 (按住shift可多选):").grid(row=2, column=0, padx=5, pady=10, sticky=tk.NW)
        self.jar_listbox = tk.Listbox(main_frame, width=60, height=8, selectmode=tk.EXTENDED)
        self.jar_listbox.grid(row=2, column=1, padx=5, pady=10)
        
        # 列表滚动条
        list_scrollbar = ttk.Scrollbar(main_frame, orient="vertical", command=self.jar_listbox.yview)
        list_scrollbar.grid(row=2, column=2, pady=10, sticky=tk.NS)
        self.jar_listbox.config(yscrollcommand=list_scrollbar.set)
        
        # 排序按钮
        ttk.Button(main_frame, text="按名称排序", command=self._sort_jars).grid(row=3, column=1, pady=5)
        
        # 状态标签
        self.status_var = tk.StringVar(value="正在搜索mods文件夹...")
        self.status_label = ttk.Label(main_frame, textvariable=self.status_var, foreground="#666")
        self.status_label.grid(row=4, column=0, columnspan=3, padx=5, pady=5, sticky=tk.W)
        
        # 处理按钮
        self.process_btn = ttk.Button(main_frame, text="开始合并", command=self._process, state=tk.DISABLED)
        self.process_btn.grid(row=5, column=1, pady=10)
        
        # 日志区域
        ttk.Label(main_frame, text="处理日志:").grid(row=6, column=0, padx=5, pady=5, sticky=tk.NW)
        self.log_text = tk.Text(main_frame, width=75, height=10, state="disabled")
        self.log_text.grid(row=6, column=1, padx=5, pady=5)
        
        # 日志滚动条
        scrollbar = ttk.Scrollbar(main_frame, orient="vertical", command=self.log_text.yview)
        scrollbar.grid(row=6, column=2, pady=5, sticky=tk.NS)
        self.log_text.config(yscrollcommand=scrollbar.set)
    
    def _log(self, message):
        """添加日志信息"""
        self.log_text.config(state="normal")
        self.log_text.insert(tk.END, message + "\n")
        self.log_text.see(tk.END)
        self.log_text.config(state="disabled")
        self.root.update_idletasks()
    
    def _auto_search_mods(self):
        """自动搜索上级目录中的mods文件夹"""
        self._log("开始自动搜索mods文件夹...")
        self._log("Fuel Loader 0.0.1-build.1加载中......")
        self._log("Fuel Loader 0.0.1-build.1加载中......")
        self._log("Fuel Loader 0.0.1-build.1加载中......")
        self._log("Fuel Loader 0.0.1-build.1加载中......")
        self._log("Fuel Loader 0.0.1-build.1加载中......")
        self._log("Fuel Loader 0.0.1-build.1加载中......")
        
        # 从当前目录向上级搜索mods文件夹（最多搜索5层）
        current_dir = os.path.abspath(os.getcwd())
        max_levels = 5
        levels_searched = 0
        
        while levels_searched < max_levels:
            candidate_mods = os.path.join(current_dir, "mods")
            if os.path.isdir(candidate_mods):
                self.mods_dir = candidate_mods
                self.mods_dir_var.set(candidate_mods)
                self._log(f"自动找到mods文件夹: {candidate_mods}")
                self._log("Fuel Loader初始化成功！")
                self._load_jar_files()
                return
            
            parent_dir = os.path.dirname(current_dir)
            if parent_dir == current_dir:  # 到达根目录
                break
            current_dir = parent_dir
            levels_searched += 1
        
        # 未找到mods文件夹
        self._log("未自动找到mods文件夹，请手动选择")
        self.status_var.set("未找到mods文件夹，请手动选择")
        self.status_label.config(foreground="#cc0000")
    
    def _select_mods_dir(self):
        """手动选择mods文件夹"""
        path = filedialog.askdirectory(title="选择mods文件夹")
        if path:
            self.mods_dir = path
            self.mods_dir_var.set(path)
            self._log(f"手动选择mods文件夹: {path}")
            self._load_jar_files()
    
    def _load_jar_files(self):
        """加载mods文件夹中的所有JAR文件"""
        if not self.mods_dir or not os.path.isdir(self.mods_dir):
            return
            
        # 清空列表
        self.jar_listbox.delete(0, tk.END)
        
        # 查找所有JAR文件
        jar_files = []
        for file in os.listdir(self.mods_dir):
            if file.lower().endswith(".jar"):
                jar_files.append(file)
        
        if not jar_files:
            self._log("mods文件夹中未找到JAR文件")
            self.status_var.set("mods文件夹中未找到JAR文件")
            self.status_label.config(foreground="#cc0000")
            return
        
        # 按名称排序并添加到列表
        jar_files.sort()
        for jar in jar_files:
            self.jar_listbox.insert(tk.END, jar)
        
        self._log(f"找到 {len(jar_files)} 个JAR文件并按名称排序")
        self.status_var.set(f"找到 {len(jar_files)} 个JAR文件，请选择目标PPTX和要合并的JAR文件")
        self._check_ready()
    
    def _sort_jars(self):
        """按名称排序JAR文件列表"""
        if self.jar_listbox.size() == 0:
            return
            
        # 获取当前列表并排序
        jars = [self.jar_listbox.get(i) for i in range(self.jar_listbox.size())]
        jars.sort()
        
        # 重新填充列表
        self.jar_listbox.delete(0, tk.END)
        for jar in jars:
            self.jar_listbox.insert(tk.END, jar)
        
        self._log("JAR文件已按名称重新排序")
    
    def _select_target_pptx(self):
        """选择目标PPTX文件"""
        path = filedialog.askopenfilename(
            title="选择要插入幻灯片的目标PPTX文件",
            filetypes=[("PPTX文件", "*.pptx"), ("所有文件", "*.*")]
        )
        if path:
            self.target_pptx_path = path
            self.target_pptx_var.set(path)
            self._log(f"已选择目标PPTX文件: {os.path.basename(path)}")
            self._check_ready()
    
    def _check_ready(self):
        """检查是否已准备好处理"""
        selected_indices = self.jar_listbox.curselection()
        if selected_indices and self.target_pptx_path:
            self.process_btn.config(state=tk.NORMAL)
            self.status_var.set(f"已选择 {len(selected_indices)} 个JAR文件，准备就绪")
            self.status_label.config(foreground="#009933")
        else:
            self.process_btn.config(state=tk.DISABLED)
    
    def _get_selected_jars(self):
        """获取选中的JAR文件列表（按名称排序）"""
        selected_indices = self.jar_listbox.curselection()
        if not selected_indices:
            return []
            
        # 按索引获取选中的JAR并保持排序
        selected_jars = [self.jar_listbox.get(i) for i in selected_indices]
        selected_jars.sort()  # 确保按名称排序
        
        # 转换为完整路径
        return [os.path.join(self.mods_dir, jar) for jar in selected_jars]
    
    def _extract_jar_contents(self, jar_path):
        """提取单个JAR文件内容到临时目录"""
        jar_name = os.path.basename(jar_path)
        self._log(f"正在提取JAR文件内容: {jar_name}")
        
        # 为每个JAR创建单独的临时目录
        jar_temp_dir = os.path.join(self.temp_dir, jar_name)
        os.makedirs(jar_temp_dir, exist_ok=True)
        
        # 提取JAR内容
        with zipfile.ZipFile(jar_path, 'r') as jar:
            jar.extractall(jar_temp_dir)
        
        return jar_temp_dir
    
    def _find_files_in_jar(self, jar_temp_dir):
        """在JAR提取目录中查找PPTX和fuel.json文件"""
        pptx_path = None
        fuel_json_path = None
        
        # 遍历目录查找文件
        for root, _, files in os.walk(jar_temp_dir):
            for file in files:
                if file.lower().endswith(".pptx") and not pptx_path:
                    pptx_path = os.path.join(root, file)
                elif file == "fuel.json" and not fuel_json_path:
                    fuel_json_path = os.path.join(root, file)
                
                if pptx_path and fuel_json_path:
                    break
        
        if not pptx_path:
            raise Exception(f"在JAR文件中未找到PPTX文件")
        if not fuel_json_path:
            raise Exception(f"在JAR文件中未找到fuel.json文件")
            
        return pptx_path, fuel_json_path
    
    def _get_insert_position(self, fuel_json_path):
        """从fuel.json获取插入位置（第二行内容）"""
        with open(fuel_json_path, 'r', encoding='utf-8') as f:
            lines = f.readlines()
        
        if len(lines) < 2:
            self._log("fuel.json文件行数不足2行，默认插入到末尾")
            return None
            
        # 获取第二行内容并尝试转换为整数位置
        position_line = lines[1].strip()
        try:
            position_str = ''.join(filter(str.isdigit, position_line))
            if not position_str:
                raise ValueError("第二行不包含数字")
                
            return int(position_str)
        except ValueError:
            self._log(f"无法将fuel.json第二行内容 '{position_line}' 转换为数字，默认插入到末尾")
            return None
    
    def _process(self):
        """执行完整处理流程"""
        try:
            # 获取选中的JAR文件
            self.selected_jars = self._get_selected_jars()
            if not self.selected_jars:
                messagebox.showwarning("警告", "请至少选择一个JAR文件")
                return
                
            self.status_var.set(f"正在处理 {len(self.selected_jars)} 个JAR文件...")
            self.status_label.config(foreground="#0066cc")
            self.process_btn.config(state=tk.DISABLED)
            self._log(f"开始处理 {len(self.selected_jars)} 个JAR文件（按名称排序）")
            
            # 收集所有JAR的处理信息
            jar_process_info = []
            for jar_path in self.selected_jars:
                jar_name = os.path.basename(jar_path)
                self._log(f"\n----- 处理 {jar_name} -----")
                
                # 提取JAR内容
                jar_temp_dir = self._extract_jar_contents(jar_path)
                
                # 查找PPTX和fuel.json
                pptx_path, fuel_json_path = self._find_files_in_jar(jar_temp_dir)
                
                # 获取插入位置
                position = self._get_insert_position(fuel_json_path)
                
                # 记录幻灯片数量
                source_prs = Presentation(pptx_path)
                slide_count = len(source_prs.slides)
                self._log(f"包含 {slide_count} 张幻灯片，插入位置: {position if position is not None else '末尾'}")
                
                # 保存处理信息
                jar_process_info.append({
                    "path": jar_path,
                    "name": jar_name,
                    "pptx": pptx_path,
                    "position": position if position is not None else float('inf'),  # 用inf表示末尾
                    "slide_count": slide_count
                })
            
            # 按位置排序，相同位置按名称排序
            jar_process_info.sort(key=lambda x: (x["position"], x["name"]))
            
            # 打开目标PPT
            target_prs = Presentation(self.target_pptx_path)
            initial_slide_count = len(target_prs.slides)
            self._log(f"\n目标PPT初始幻灯片数量: {initial_slide_count}")
            
            # 跟踪位置偏移量（因插入幻灯片导致的位置变化）
            position_offset = 0
            
            # 逐个处理JAR并插入幻灯片
            for info in jar_process_info:
                self._log(f"\n----- 插入 {info['name']} -----")
                
                # 计算实际插入位置
                actual_position = info["position"] + position_offset
                # 确保位置在有效范围内
                if actual_position < 0 or actual_position > len(target_prs.slides):
                    actual_position = len(target_prs.slides)
                    self._log(f"调整插入位置为: {actual_position} (末尾)")
                
                # 打开源PPT
                source_prs = Presentation(info["pptx"])
                
                # 插入幻灯片
                for i, slide in enumerate(source_prs.slides):
                    # 计算当前插入位置
                    current_position = actual_position + i
                    
                    # 复制幻灯片
                    slide_layout = target_prs.slide_layouts[6]  # 空白布局
                    new_slide = target_prs.slides.add_slide(slide_layout)
                    
                    # 复制内容
                    self._copy_slide_content(slide, new_slide)
                    self._log(f"已插入第 {i+1}/{info['slide_count']} 张幻灯片 (位置: {current_position})")
                
                # 更新位置偏移量
                position_offset += info["slide_count"]
                self._log(f"{info['name']} 插入完成，新增 {info['slide_count']} 张幻灯片")
            
            # 保存合并后的PPT
            output_path = os.path.splitext(self.target_pptx_path)[0] + "_merged.pptx"
            target_prs.save(output_path)
            
            # 处理完成
            final_slide_count = len(target_prs.slides)
            self._log(f"\n所有JAR文件处理完成！")
            self._log(f"原始幻灯片: {initial_slide_count} 张")
            self._log(f"新增幻灯片: {final_slide_count - initial_slide_count} 张")
            self._log(f"最终幻灯片: {final_slide_count} 张")
            self._log(f"文件已保存至: {output_path}")
            
            self.status_var.set("所有JAR文件合并完成！")
            self.status_label.config(foreground="#009933")
            messagebox.showinfo("成功", f"PPT合并完成！\n共处理 {len(self.selected_jars)} 个JAR文件\n新增 {final_slide_count - initial_slide_count} 张幻灯片\n文件已保存至:\n{output_path}")
            
        except Exception as e:
            error_msg = f"处理失败: {str(e)}"
            self._log(error_msg)
            self.status_var.set(error_msg)
            self.status_label.config(foreground="#cc0000")
            messagebox.showerror("错误", error_msg)
        finally:
            self.process_btn.config(state=tk.NORMAL)
    
    def _copy_slide_content(self, source_slide, target_slide):
        """复制幻灯片内容（文本框和图片）"""
        for shape in source_slide.shapes:
            # 处理文本框
            if shape.has_text_frame:
                new_shape = target_slide.shapes.add_textbox(
                    shape.left, shape.top, shape.width, shape.height
                )
                new_tf = new_shape.text_frame
                # 复制文本内容
                for para in shape.text_frame.paragraphs:
                    new_para = new_tf.add_paragraph()
                    new_para.text = para.text
                    new_para.font.name = para.font.name
                    new_para.font.size = para.font.size
                    new_para.font.bold = para.font.bold
                    new_para.font.italic = para.font.italic
                    new_para.alignment = para.alignment
                    # 安全处理颜色属性
                    if hasattr(para.font.color, 'rgb') and para.font.color.rgb is not None:
                        new_para.font.color.rgb = para.font.color.rgb
            
            # 处理图片
            elif shape.shape_type == 13:  # 13是图片类型
                image_bytes = shape.image.blob
                target_slide.shapes.add_picture(
                    BytesIO(image_bytes),
                    shape.left, shape.top,
                    shape.width, shape.height
                )
    
    def __del__(self):
        """清理临时目录"""
        if os.path.exists(self.temp_dir):
            shutil.rmtree(self.temp_dir, ignore_errors=True)

if __name__ == "__main__":
    # 检查依赖
    try:
        import pptx
    except ImportError:
        print("请先安装依赖: pip install python-pptx")
        sys.exit(1)
    
    root = tk.Tk()
    app = MultiJarPPTMerger(root)
    root.mainloop()
    
